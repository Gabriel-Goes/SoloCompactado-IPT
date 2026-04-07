#!/usr/bin/env python3
"""
Gera a apresentação 'Mapeamento de Análogos ao Terranimo' usando template IPT.
Uso: python scripts/build_slides.py
Saída: docs/apresentacao_analogos_terranimo_2026-04-05.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# ── Caminhos ──────────────────────────────────────────────────────────────────
ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "docs" / "Proposta_Compactação_Solos.pptx"
OUTPUT   = ROOT / "docs" / "apresentacao_analogos_terranimo_2026-04-05.pptx"
SCREENSHOTS = ROOT / "terranimo" / "screenshots"
IMG_WEB  = ROOT / "docs" / "img_web"
IMG_DL   = ROOT / "docs" / "img_slides"

# ── Cores do template ─────────────────────────────────────────────────────────
VERDE    = RGBColor(0x4C, 0xAF, 0x50)
AMARELO  = RGBColor(0xFF, 0xC1, 0x07)
VERMELHO = RGBColor(0xF4, 0x43, 0x36)
CINZA_ESC = RGBColor(0x33, 0x33, 0x33)
CINZA_MED = RGBColor(0x66, 0x66, 0x66)
BRANCO   = RGBColor(0xFF, 0xFF, 0xFF)
AZUL_IPT = RGBColor(0x00, 0x3D, 0x6B)

# ── Helpers ───────────────────────────────────────────────────────────────────

def set_text(placeholder, text, font_size=None, bold=False, color=None, alignment=None):
    """Define texto em um placeholder, limpando conteúdo anterior."""
    tf = placeholder.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    if font_size:
        run.font.size = Pt(font_size)
    if bold:
        run.font.bold = True
    if color:
        run.font.color.rgb = color
    if alignment:
        p.alignment = alignment


def add_bullet_slide(prs, layout_idx, title, bullets, font_size=16):
    """Slide com título e lista de bullets."""
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    # Título
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, title, font_size=28, bold=True)
            break
    # Conteúdo
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 1:
            tf = ph.text_frame
            tf.clear()
            for i, bullet in enumerate(bullets):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                run = p.add_run()
                run.text = bullet
                run.font.size = Pt(font_size)
                run.font.color.rgb = CINZA_ESC
                p.space_after = Pt(6)
            break
    return slide


def add_section_slide(prs, title):
    """Slide separador de seção (layout 17)."""
    slide = prs.slides.add_slide(prs.slide_layouts[17])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 14:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = title
            run.font.size = Pt(36)
            run.font.bold = True
            break
    return slide


def add_image_to_slide(slide, img_path, left, top, width=None, height=None):
    """Adiciona imagem ao slide."""
    img_path = str(img_path)
    if width and height:
        slide.shapes.add_picture(img_path, left, top, width, height)
    elif width:
        slide.shapes.add_picture(img_path, left, top, width=width)
    elif height:
        slide.shapes.add_picture(img_path, left, top, height=height)
    else:
        slide.shapes.add_picture(img_path, left, top)


def add_textbox(slide, left, top, width, height, text, font_size=14, bold=False,
                color=CINZA_ESC, alignment=PP_ALIGN.LEFT, word_wrap=True):
    """Adiciona caixa de texto ao slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    p.alignment = alignment
    return txBox


def add_semaforo_box(slide, left, top, width, height, cor, label):
    """Adiciona retângulo colorido do semáforo."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = cor
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = BRANCO if cor != AMARELO else CINZA_ESC


# ── Construção da apresentação ────────────────────────────────────────────────

def build():
    prs = Presentation(str(TEMPLATE))

    # Remove slides existentes do template (5 slides do Proposta)
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId is None:
            # Try alternative: get the relationship id
            slide_id_elem = prs.slides._sldIdLst[0]
            rId = slide_id_elem.attrib.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
        if rId:
            prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

    # ╔═══════════════════════════════════════════════════════════════════════╗
    # ║  SLIDE 1 — CAPA                                                      ║
    # ╚═══════════════════════════════════════════════════════════════════════╝
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # Título Vertical
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:  # Título
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Compactação de Solos\npor Tráfego de Máquinas"
            run.font.size = Pt(32)
            run.font.bold = True
        elif idx == 1:  # Subtítulo
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "IPT — CIMA"
            run.font.size = Pt(18)
            p2 = tf.add_paragraph()
            run2 = p2.add_run()
            run2.text = "Mapeamento de soluções e proposta para o Brasil"
            run2.font.size = Pt(14)
            p3 = tf.add_paragraph()
            run3 = p3.add_run()
            run3.text = "Abril 2026"
            run3.font.size = Pt(12)

    # ╔═══════════════════════════════════════════════════════════════════════╗
    # ║  SLIDE 2 — O PROBLEMA                                                ║
    # ╚═══════════════════════════════════════════════════════════════════════╝
    add_section_slide(prs, "O Problema")

    # ╔═══════════════════════════════════════════════════════════════════════╗
    # ║  SLIDE 3 — Impacto da compactação                                    ║
    # ╚═══════════════════════════════════════════════════════════════════════╝
    slide = prs.slides.add_slide(prs.slide_layouts[2])  # Título e Conteúdo
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            set_text(ph, "Compactação: um problema silencioso", font_size=28, bold=True)
        elif idx == 1:
            tf = ph.text_frame
            tf.clear()
            dados = [
                ("Perda de produtividade:", "6–34% de queda na produção de grãos (meta-análise global)"),
                ("Raízes:", "crescimento freia a 100 kPa e para a 1000 kPa de resistência"),
                ("Peso das máquinas:", "rodas de colhedoras passaram de 1,5 t (1960) para 9 t (hoje)"),
                ("Irreversibilidade:", "compactação de subsuperfície persiste por décadas"),
                ("Cana-de-açúcar:", "fator mais restritivo de produção há 40+ anos"),
                ("", ""),
                ("Fazendeiro → CNH → IPT:", "\"Como evitar compactação em tempo real?\""),
            ]
            for i, (titulo_b, desc) in enumerate(dados):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                if titulo_b == "":
                    p.space_before = Pt(12)
                    continue
                if titulo_b.startswith("Fazendeiro"):
                    run = p.add_run()
                    run.text = titulo_b + " " + desc
                    run.font.size = Pt(18)
                    run.font.bold = True
                    run.font.color.rgb = AZUL_IPT
                    p.space_before = Pt(12)
                else:
                    run1 = p.add_run()
                    run1.text = titulo_b + " "
                    run1.font.size = Pt(16)
                    run1.font.bold = True
                    run1.font.color.rgb = CINZA_ESC
                    run2 = p.add_run()
                    run2.text = desc
                    run2.font.size = Pt(16)
                    run2.font.color.rgb = CINZA_MED
                    p.space_after = Pt(4)

    # ╔═══════════════════════════════════════════════════════════════════════╗
    # ║  SLIDE 4 — SEÇÃO: Soluções no mundo                                  ║
    # ╚═══════════════════════════════════════════════════════════════════════╝
    add_section_slide(prs, "Soluções no Mundo")

    # ╔═══════════════════════════════════════════════════════════════════════╗
    # ║  SLIDE 5 — Terranimo: a referência (com screenshots)                  ║
    # ╚═══════════════════════════════════════════════════════════════════════╝
    slide = prs.slides.add_slide(prs.slide_layouts[18])  # Somente Título
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "Terranimo — a referência mundial", font_size=28, bold=True)

    # Texto à esquerda
    bullets_text = (
        "Consórcio: Agroscope (CH), Aarhus (DK),\n"
        "SLU (SE), Wageningen (NL)\n\n"
        "Modo Light: 4 entradas → semáforo\n"
        "  • Carga na roda\n"
        "  • Pressão do pneu\n"
        "  • Tipo de solo\n"
        "  • Umidade\n\n"
        "Modo Expert: perfil por camadas,\n"
        "parâmetros mecânicos completos\n\n"
        "Banco de máquinas pré-configurado\n"
        "Interface web gratuita (terranimo.ch)"
    )
    add_textbox(slide, Cm(0.8), Cm(4.5), Cm(10), Cm(12),
                bullets_text, font_size=13, color=CINZA_ESC)

    # Screenshots do Terranimo à direita (2x2 grid)
    imgs = [
        SCREENSHOTS / "jc-20260220-210212.png",  # máquina + solo
        SCREENSHOTS / "jc-20260220-210206.png",  # pneus + pressão
        SCREENSHOTS / "jc-20260220-210217.png",  # tipo de solo
        SCREENSHOTS / "jc-20260220-210231.png",  # resultado semáforo
    ]
    positions = [
        (Cm(12.5), Cm(3.8)),  # top-left
        (Cm(21.0), Cm(3.8)),  # top-right
        (Cm(12.5), Cm(9.8)),  # bottom-left
        (Cm(21.0), Cm(9.8)),  # bottom-right
    ]
    img_w = Cm(8.0)
    img_h = Cm(5.5)
    for img_path, (x, y) in zip(imgs, positions):
        if img_path.exists():
            add_image_to_slide(slide, img_path, x, y, width=img_w, height=img_h)

    # ╔═══════════════════════════════════════════════════════════════════════╗
    # ║  SLIDE 6 — Terranimo: o semáforo                                      ║
    # ╚═══════════════════════════════════════════════════════════════════════╝
    slide = prs.slides.add_slide(prs.slide_layouts[18])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "Como funciona: o semáforo de risco", font_size=28, bold=True)

    # Screenshot grande do resultado
    result_img = SCREENSHOTS / "jc-20260220-210231.png"
    if result_img.exists():
        add_image_to_slide(slide, result_img, Cm(0.8), Cm(3.5), width=Cm(16), height=Cm(11))

    # Semáforos à direita
    x_sem = Cm(18.5)
    w_sem = Cm(12.5)
    h_sem = Cm(2.8)

    add_semaforo_box(slide, x_sem, Cm(3.5), w_sem, h_sem, VERDE,
                     "VERDE\nσ < 0,5 × σp\nSem risco — pode trafegar")
    add_semaforo_box(slide, x_sem, Cm(6.8), w_sem, h_sem, AMARELO,
                     "AMARELO\n0,5 × σp ≤ σ < 1,1 × σp\nAtenção — risco moderado")
    add_semaforo_box(slide, x_sem, Cm(10.1), w_sem, h_sem, VERMELHO,
                     "VERMELHO\nσ ≥ 1,1 × σp\nCompactação — alterar pressão/rota")

    add_textbox(slide, x_sem, Cm(13.5), w_sem, Cm(2),
                "σ = tensão aplicada pela máquina\nσp = resistência do solo (pré-adensamento)",
                font_size=11, color=CINZA_MED)

    # ╔═══════════════════════════════════════════════════════════════════════╗
    # ║  SLIDE 7 — CLAAS CEMOS + Terranimo embarcado                          ║
    # ╚═══════════════════════════════════════════════════════════════════════╝
    slide = prs.slides.add_slide(prs.slide_layouts[18])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "CLAAS CEMOS — Terranimo no trator", font_size=28, bold=True)

    # Imagem CLAAS grande
    cemos_img = IMG_WEB / "Claas-integrates-Terranimo-into-Cemos-for-tractors.jpg"
    if cemos_img.exists():
        add_image_to_slide(slide, cemos_img, Cm(1), Cm(3.5), width=Cm(17), height=Cm(11.5))

    # Texto à direita
    cemos_text = (
        "Desde 2022: Terranimo integrado\n"
        "ao terminal CEBIS do trator CLAAS\n\n"
        "O que faz:\n"
        "• Avalia risco de compactação\n"
        "  em tempo real\n"
        "• Recomenda ajuste de pressão\n"
        "  dos pneus\n"
        "• Recomenda lastragem ideal\n\n"
        "Medalha de prata Agritechnica 2022\n\n"
        "Referência para o que queremos\n"
        "construir para a CNH no Brasil"
    )
    add_textbox(slide, Cm(19), Cm(3.5), Cm(12), Cm(12),
                cemos_text, font_size=14, color=CINZA_ESC)

    # ╔═══════════════════════════════════════════════════════════════════════╗
    # ║  SLIDE 8 — Soil2Cover: a fronteira                                    ║
    # ╚═══════════════════════════════════════════════════════════════════════╝
    slide = prs.slides.add_slide(prs.slide_layouts[18])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "Soil2Cover — otimização de rota (fronteira)", font_size=28, bold=True)

    # Flowchart à esquerda
    fig1 = IMG_DL / "soil2cover_fig1_flowchart.jpg"
    if fig1.exists():
        add_image_to_slide(slide, fig1, Cm(0.5), Cm(3.2), height=Cm(13))

    # Route map ao centro
    fig7 = IMG_DL / "soil2cover_fig7_route.jpg"
    if fig7.exists():
        add_image_to_slide(slide, fig7, Cm(11), Cm(3.5), width=Cm(11), height=Cm(5.5))

    # Texto à direita
    s2c_text = (
        "Keller et al. (2025)\n"
        "Precision Agriculture\n\n"
        "Não só avalia risco:\n"
        "otimiza a rota que minimiza\n"
        "compactação acumulada\n\n"
        "• SoilFlex como função-objetivo\n"
        "• Reduz rota em 4–6%\n"
        "• Reduz compactação em\n"
        "  cabeceiras em até 30%\n"
        "• Testado em 1000+ campos\n\n"
        "É o objetivo final do nosso\n"
        "projeto para a CNH"
    )
    add_textbox(slide, Cm(11), Cm(9.5), Cm(12), Cm(8),
                s2c_text, font_size=13, color=CINZA_ESC)

    # Ref
    add_textbox(slide, Cm(23), Cm(3.5), Cm(9), Cm(2),
                "Mier et al., Precision Agriculture 26, 57 (2025)",
                font_size=9, color=CINZA_MED)

    # ╔═══════════════════════════════════════════════════════════════════════╗
    # ║  SLIDE 9 — Quadro comparativo internacional                           ║
    # ╚═══════════════════════════════════════════════════════════════════════╝
    slide = prs.slides.add_slide(prs.slide_layouts[18])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "Quadro comparativo — ferramentas no mundo", font_size=28, bold=True)

    # Tabela
    rows, cols = 8, 7
    tbl_w = Cm(31)
    tbl_h = Cm(12)
    table_shape = slide.shapes.add_table(rows, cols, Cm(1), Cm(3.5), tbl_w, tbl_h)
    table = table_shape.table

    headers = ["Ferramenta", "País", "Motor", "Embarcável", "Otim. Rota", "PTF Brasil", "Open Source"]
    data = [
        ["Terranimo",    "CH/DK/SE/NL", "SoilFlex",   "✓ (CLAAS)", "—", "—",  "—"],
        ["Soil2Cover",   "NL/ES",       "SoilFlex",   "—",         "✓", "—",  "—"],
        ["SoilFlex",     "CH",          "Analítico",  "—",         "—", "—",  "—"],
        ["SOCOMO",       "NL",          "Boussinesq", "—",         "—", "—",  "—"],
        ["REPRO",        "DE",          "Vários",     "—",         "—", "—",  "—"],
        ["soilphysics",  "BR",          "SoilFlex",   "—",         "—", "✓",  "✓"],
        ["COMPSOHMA",    "BR",          "SoilFlex",   "—",         "—", "✓*", "—"],
    ]

    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(12)
                run.font.bold = True
                run.font.color.rgb = BRANCO
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZUL_IPT

    for i, row_data in enumerate(data):
        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(11)
                    run.font.color.rgb = CINZA_ESC
            # Highlight Brazil rows
            if row_data[0] in ("soilphysics", "COMPSOHMA"):
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)

    # Nota de rodapé
    add_textbox(slide, Cm(1), Cm(16), Cm(20), Cm(1),
                "* COMPSOHMA: PTFs limitadas a cana-de-açúcar. soilphysics: PTFs para Latossolos (Severiano et al. 2013).",
                font_size=10, color=CINZA_MED)

    # ╔═══════════════════════════════════════════════════════════════════════╗
    # ║  SLIDE 10 — SEÇÃO: E no Brasil?                                       ║
    # ╚═══════════════════════════════════════════════════════════════════════╝
    add_section_slide(prs, "E no Brasil?")

    # ╔═══════════════════════════════════════════════════════════════════════╗
    # ║  SLIDE 11 — soilphysics: a base científica                            ║
    # ╚═══════════════════════════════════════════════════════════════════════╝
    slide = prs.slides.add_slide(prs.slide_layouts[18])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "soilphysics — motor brasileiro open source", font_size=28, bold=True)

    # Imagem contact area
    ca_img = IMG_DL / "soilphysics_contactarea.jpeg"
    if ca_img.exists():
        add_image_to_slide(slide, ca_img, Cm(1), Cm(3.5), height=Cm(8))

    # Imagem fluxograma sigmaP
    sp_img = IMG_DL / "soilphysics_sigmap.jpg"
    if sp_img.exists():
        add_image_to_slide(slide, sp_img, Cm(9), Cm(3.5), height=Cm(10))

    # Texto à direita
    sp_text = (
        "Pacote R — ESALQ/USP + IF Goiano\n"
        "(de Lima, da Silva, Tormena et al.)\n\n"
        "3 funções-chave:\n"
        "• stressTraffic() → tensão por camada\n"
        "  (TASC + Boussinesq/Froehlich)\n"
        "• soilStrength() → resistência σp\n"
        "  (PTFs brasileiras: Severiano, Imhoff)\n"
        "• soilDeformation() → Δρ previsto\n"
        "  (O'Sullivan & Robertson 1996)\n\n"
        "Mesmo motor do Terranimo\n"
        "Código aberto no CRAN\n"
        "Registro INPI"
    )
    add_textbox(slide, Cm(19), Cm(3.2), Cm(13), Cm(13),
                sp_text, font_size=13, color=CINZA_ESC)

    # ╔═══════════════════════════════════════════════════════════════════════╗
    # ║  SLIDE 12 — PredComp vs COMPSOHMA                                     ║
    # ╚═══════════════════════════════════════════════════════════════════════╝
    slide = prs.slides.add_slide(prs.slide_layouts[18])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "Ferramentas brasileiras: PredComp e COMPSOHMA", font_size=26, bold=True)

    # Duas colunas com boxes
    # PredComp (esquerda)
    box_w = Cm(14)
    box_h = Cm(11)

    shape_l = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), Cm(3.5), box_w, box_h)
    shape_l.fill.solid()
    shape_l.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    shape_l.line.color.rgb = AZUL_IPT
    tf = shape_l.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "PredComp (acadêmico)"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = AZUL_IPT

    pred_items = [
        "Base: soilphysics / SoilFlex",
        "5 abas: tensão, σp, compressão, risco, predição",
        "Entradas detalhadas: pressão, carga, N, λ, κ, ρ",
        "Exige conhecimento de parâmetros mecânicos",
        "Ativo no CRAN + registro INPI",
        "FAPESP + CNPq",
    ]
    for item in pred_items:
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = "• " + item
        run.font.size = Pt(12)
        run.font.color.rgb = CINZA_ESC

    # COMPSOHMA (direita)
    shape_r = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(17), Cm(3.5), box_w, box_h)
    shape_r.fill.solid()
    shape_r.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
    shape_r.line.color.rgb = VERDE
    tf = shape_r.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "COMPSOHMA (aplicado, beta)"
    run.font.size = Pt(18)
    run.font.bold = True
    run.font.color.rgb = VERDE

    comp_items = [
        "Estilo Terranimo Light",
        "Risco até 1 m (4 camadas)",
        "Cenários pré-configurados (cana)",
        "Entradas simples: máquina, argila, umidade",
        "Saída: perfil verde / amarelo / vermelho",
        "Beta — foco atual: cana-de-açúcar",
    ]
    for item in comp_items:
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = "• " + item
        run.font.size = Pt(12)
        run.font.color.rgb = CINZA_ESC

    # Nota
    add_textbox(slide, Cm(1), Cm(15.2), Cm(30), Cm(1.5),
                "Ambos desenvolvidos na ESALQ/USP (Piracicaba). Nenhum é embarcável ou integra GPS.",
                font_size=12, bold=True, color=VERMELHO)

    # ╔═══════════════════════════════════════════════════════════════════════╗
    # ║  SLIDE 13 — PTFs brasileiras                                          ║
    # ╚═══════════════════════════════════════════════════════════════════════╝
    slide = prs.slides.add_slide(prs.slide_layouts[18])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "PTFs calibradas para solos brasileiros", font_size=28, bold=True)

    # Tabela de PTFs
    rows, cols = 6, 4
    tbl = slide.shapes.add_table(rows, cols, Cm(1), Cm(3.5), Cm(22), Cm(8))
    table = tbl.table

    ptf_headers = ["Referência", "Variável", "Solos", "Cobertura"]
    ptf_data = [
        ["Severiano et al. (2013)", "σp = A × h^B", "Latossolos (152–521 g/kg argila)", "✓ Boa"],
        ["Imhoff et al. (2004)", "σp = f(ρ, CC, w)", "Hapludox", "✓ Boa"],
        ["de Lima et al. (2018)", "N, λ, κ", "Franco-arenoso / Franco-argiloarenoso", "Parcial"],
        ["de Lima et al. (2020)", "Camada arável e pé-de-grade", "Solos agrícolas SP", "Parcial"],
        ["Horn & Fleige (2003)", "σp europeu", "Solos europeus (referência)", "Não tropical"],
    ]

    for j, h in enumerate(ptf_headers):
        cell = table.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            for run in paragraph.runs:
                run.font.size = Pt(12)
                run.font.bold = True
                run.font.color.rgb = BRANCO
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZUL_IPT

    for i, row_data in enumerate(ptf_data):
        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = val
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.size = Pt(11)

    # Lacunas à direita
    lacuna_text = (
        "Lacunas de cobertura:\n\n"
        "• Argissolos\n"
        "• Nitossolos\n"
        "• Cambissolos\n"
        "• Neossolos\n\n"
        "Necessário: campanha de\n"
        "calibração com novos solos"
    )
    add_textbox(slide, Cm(24.5), Cm(3.5), Cm(8), Cm(10),
                lacuna_text, font_size=13, color=CINZA_ESC)

    # ╔═══════════════════════════════════════════════════════════════════════╗
    # ║  SLIDE 14 — Lacunas críticas (visual)                                 ║
    # ╚═══════════════════════════════════════════════════════════════════════╝
    slide = prs.slides.add_slide(prs.slide_layouts[18])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "O que falta para um Terranimo brasileiro?", font_size=28, bold=True)

    lacunas = [
        ("Sem ferramenta agnóstica de cultura", "COMPSOHMA só cana", VERMELHO),
        ("Sem integração espacial GPS/RTK", "Nada gera mapa de risco por trilha", VERMELHO),
        ("Sem capacidade embarcada", "Nada como CLAAS CEMOS no Brasil", VERMELHO),
        ("Cobertura incompleta de PTFs", "Faltam Argissolos, Nitossolos, Cambissolos", AMARELO),
        ("Sem conexão com BDSolos/SiBCS", "Entrada toda manual", AMARELO),
        ("Sem framework regulatório", "Brasil não tem equivalente ao BBodSchG alemão", AMARELO),
    ]

    y_start = Cm(3.8)
    for i, (titulo, desc, cor) in enumerate(lacunas):
        y = y_start + Cm(2.0) * i
        # Barra colorida
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(1), y, Cm(0.6), Cm(1.5))
        bar.fill.solid()
        bar.fill.fore_color.rgb = cor
        bar.line.fill.background()

        # Título
        add_textbox(slide, Cm(2.2), y, Cm(15), Cm(1),
                    titulo, font_size=15, bold=True, color=CINZA_ESC)
        # Descrição
        add_textbox(slide, Cm(2.2), y + Cm(0.8), Cm(15), Cm(1),
                    desc, font_size=12, color=CINZA_MED)

    # Conclusão à direita
    concl_text = (
        "A ciência existe.\n"
        "O motor existe (soilphysics).\n"
        "As PTFs existem (parcialmente).\n\n"
        "Falta a integração:\n"
        "motor + GPS + interface +\n"
        "banco de máquinas +\n"
        "embarcar no trator"
    )
    box_c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Cm(19), Cm(4), Cm(12.5), Cm(10))
    box_c.fill.solid()
    box_c.fill.fore_color.rgb = RGBColor(0xE3, 0xF2, 0xFD)
    box_c.line.color.rgb = AZUL_IPT
    tf = box_c.text_frame
    tf.word_wrap = True
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = concl_text
    run.font.size = Pt(16)
    run.font.color.rgb = AZUL_IPT
    run.font.bold = True

    # ╔═══════════════════════════════════════════════════════════════════════╗
    # ║  SLIDE 15 — SEÇÃO: Proposta                                           ║
    # ╚═══════════════════════════════════════════════════════════════════════╝
    add_section_slide(prs, "Proposta IPT")

    # ╔═══════════════════════════════════════════════════════════════════════╗
    # ║  SLIDE 16 — Cadeia de cálculo (diagrama)                              ║
    # ╚═══════════════════════════════════════════════════════════════════════╝
    slide = prs.slides.add_slide(prs.slide_layouts[18])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "Motor de cálculo — do pneu ao semáforo", font_size=28, bold=True)

    # Diagrama com shapes
    # Coluna 1: Entrada Máquina
    def add_box(slide, x, y, w, h, text, fill_color, text_color=BRANCO, font_sz=11):
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_sz)
        run.font.bold = True
        run.font.color.rgb = text_color
        return shape

    box_w = Cm(7)
    box_h = Cm(2.5)

    # Linha 1: Entradas
    add_box(slide, Cm(1.5), Cm(3.5), box_w, box_h,
            "MÁQUINA\nPressão, diâmetro,\nlargura, carga", RGBColor(0x42, 0x6E, 0xB4))
    add_box(slide, Cm(12.5), Cm(3.5), box_w, box_h,
            "SOLO\nArgila%, sucção,\nξ, ρ por camada", RGBColor(0x8D, 0x6E, 0x63))
    add_box(slide, Cm(23.5), Cm(3.5), box_w, box_h,
            "GPS / RTK\nRota, velocidade,\npassadas", RGBColor(0x00, 0x79, 0x6B))

    # Linha 2: Motor
    add_box(slide, Cm(1.5), Cm(7), box_w, Cm(2),
            "TASC → Área de contato", RGBColor(0x56, 0x85, 0xC7))
    add_box(slide, Cm(12.5), Cm(7), box_w, Cm(2),
            "PTF → σp por camada", RGBColor(0xA1, 0x88, 0x7F))

    # Linha 3: Propagação
    add_box(slide, Cm(7), Cm(10), Cm(10), Cm(2),
            "Boussinesq/Froehlich → σ aplicada(z)", RGBColor(0x37, 0x47, 0x4F), BRANCO, 12)

    # Linha 4: Comparação
    add_box(slide, Cm(7), Cm(12.8), Cm(10), Cm(1.8),
            "σ aplicada  vs  σp  →  SEMÁFORO", RGBColor(0x1B, 0x5E, 0x20), BRANCO, 13)

    # Semáforos
    add_semaforo_box(slide, Cm(19), Cm(10), Cm(3.5), Cm(1.2), VERDE, "VERDE")
    add_semaforo_box(slide, Cm(23), Cm(10), Cm(3.5), Cm(1.2), AMARELO, "AMARELO")
    add_semaforo_box(slide, Cm(27), Cm(10), Cm(3.5), Cm(1.2), VERMELHO, "VERMELHO")

    # Setas (usando textbox com →)
    add_textbox(slide, Cm(4), Cm(6.2), Cm(2), Cm(1), "↓", font_size=20, bold=True,
                color=CINZA_MED, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(15), Cm(6.2), Cm(2), Cm(1), "↓", font_size=20, bold=True,
                color=CINZA_MED, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(11), Cm(9.2), Cm(2), Cm(1), "↓", font_size=20, bold=True,
                color=CINZA_MED, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Cm(11), Cm(12), Cm(2), Cm(1), "↓", font_size=20, bold=True,
                color=CINZA_MED, alignment=PP_ALIGN.CENTER)

    # Mapa final
    add_box(slide, Cm(19), Cm(12.8), Cm(12), Cm(1.8),
            "MAPA DE RISCO por ponto GPS", RGBColor(0x00, 0x79, 0x6B), BRANCO, 12)

    # ╔═══════════════════════════════════════════════════════════════════════╗
    # ║  SLIDE 17 — Roadmap / Próximos passos                                ║
    # ╚═══════════════════════════════════════════════════════════════════════╝
    slide = prs.slides.add_slide(prs.slide_layouts[18])
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            set_text(ph, "Roadmap — do protótipo ao trator", font_size=28, bold=True)

    phases = [
        ("FASE 1", "Validação do motor", "Reimplementar SoilFlex em Python\nValidar contra soilphysics (R)\nCenários: Latossolo + colhedora/trator", RGBColor(0x42, 0x6E, 0xB4)),
        ("FASE 2", "Protótipo visual", "Dashboard Streamlit + PyDeck\nGêmeo digital: trator na rota\nMapa de risco em tempo real", RGBColor(0x2E, 0x7D, 0x32)),
        ("FASE 3", "Integração GPS", "Dados RTK reais\nMapa espacial de solo\nCompactação acumulada", RGBColor(0xF5, 0x7F, 0x17)),
        ("FASE 4", "Embarcar", "Otimizar para C/C++\nIntegrar ao computador do trator\nProtocolo ISOBUS (ISO 11783)", RGBColor(0xC6, 0x28, 0x28)),
    ]

    for i, (fase, titulo, desc, cor) in enumerate(phases):
        x = Cm(1) + Cm(7.8) * i
        y = Cm(4)

        # Número da fase
        fase_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Cm(7), Cm(2))
        fase_box.fill.solid()
        fase_box.fill.fore_color.rgb = cor
        fase_box.line.fill.background()
        tf = fase_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"{fase}\n{titulo}"
        run.font.size = Pt(14)
        run.font.bold = True
        run.font.color.rgb = BRANCO

        # Descrição
        desc_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Cm(6.5), Cm(7), Cm(7))
        desc_box.fill.solid()
        desc_box.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)
        desc_box.line.color.rgb = cor
        tf = desc_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = desc
        run.font.size = Pt(12)
        run.font.color.rgb = CINZA_ESC

        # Seta entre fases
        if i < 3:
            add_textbox(slide, x + Cm(7.1), Cm(4.5), Cm(1), Cm(1), "→",
                       font_size=24, bold=True, color=CINZA_MED, alignment=PP_ALIGN.CENTER)

    # Nota
    add_textbox(slide, Cm(1), Cm(14.5), Cm(30), Cm(1.5),
                "Base tecnológica: soilphysics (motor SoilFlex) + PTFs Severiano → mesmo fundamento do Terranimo + calibração tropical",
                font_size=12, bold=True, color=AZUL_IPT)

    # ╔═══════════════════════════════════════════════════════════════════════╗
    # ║  SLIDE 18 — Grupos de pesquisa                                        ║
    # ╚═══════════════════════════════════════════════════════════════════════╝
    slide = prs.slides.add_slide(prs.slide_layouts[2])
    for ph in slide.placeholders:
        idx = ph.placeholder_format.idx
        if idx == 0:
            set_text(ph, "Parceiros científicos potenciais", font_size=28, bold=True)
        elif idx == 1:
            tf = ph.text_frame
            tf.clear()
            groups = [
                "ESALQ/USP (Piracicaba) — R. P. de Lima, M. R. Cherubin, A. P. da Silva",
                "    → soilphysics, PredComp, COMPSOHMA, probe automatizada",
                "IF Goiano (Urutaí) — A. R. da Silva, co-desenvolvedor soilphysics",
                "UFLA (Lavras) — Severiano et al., PTFs para Latossolos",
                "UEM/UEPG (Paraná) — Tormena, Giarola: SoilFlex-LLWR com Keller",
                "UFRPE (Recife) — ensaios de rodagem com solos brasileiros",
                "EMBRAPA Solos — BDSolos, SiBCS, SoloClass (sem compactação)",
                "",
                "Recursos: BDSolos (EMBRAPA), SiBCS, MapBiomas Solo, FEBR",
            ]
            for i, g in enumerate(groups):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                run = p.add_run()
                run.text = g
                run.font.size = Pt(14)
                if g.startswith("    →"):
                    run.font.color.rgb = CINZA_MED
                    run.font.size = Pt(12)
                elif g.startswith("Recursos"):
                    run.font.bold = True
                    run.font.color.rgb = AZUL_IPT
                else:
                    run.font.color.rgb = CINZA_ESC
                p.space_after = Pt(3)

    # ╔═══════════════════════════════════════════════════════════════════════╗
    # ║  SLIDE 19 — Obrigado / Contatos                                       ║
    # ╚═══════════════════════════════════════════════════════════════════════╝
    slide = prs.slides.add_slide(prs.slide_layouts[20])  # Contatos
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 12:
            tf = ph.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = "Obrigado!"
            run.font.size = Pt(36)
            run.font.bold = True

            contatos = [
                "",
                "Eng. Dr. Scandar Ignatius (ignatius@ipt.br) — pesquisador",
                "Tec. Rubens Vieira (rvieira@ipt.br) — pesquisador",
                "Eng. Me. Rodrigo Serafim (rodrigos@ipt.br) — técnico",
                "Geofís. Dr. Otávio Brasil (gandolfo@ipt.br) — pesquisador",
                "Geól. Me Daniel Seabra (dseabra@ipt.br) — gerente de negócios",
            ]
            for c in contatos:
                p = tf.add_paragraph()
                run = p.add_run()
                run.text = c
                run.font.size = Pt(14)
                run.font.color.rgb = CINZA_ESC

    # ── Salvar ────────────────────────────────────────────────────────────────
    prs.save(str(OUTPUT))
    print(f"Apresentação salva em: {OUTPUT}")
    print(f"Total de slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
