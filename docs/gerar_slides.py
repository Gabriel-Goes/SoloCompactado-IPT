#!/usr/bin/env python3
"""Gera apresentação PPTX: Mapeamento de Análogos ao Terranimo."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ── Paleta ──────────────────────────────────────────────────────────
AZUL_ESCURO   = RGBColor(0x14, 0x2D, 0x4C)
AZUL_MEDIO    = RGBColor(0x1A, 0x5C, 0x8E)
AZUL_ACCENT   = RGBColor(0x0D, 0x6E, 0xAA)
AZUL_CLARO_BG = RGBColor(0xE8, 0xF0, 0xF8)
BRANCO        = RGBColor(0xFF, 0xFF, 0xFF)
CINZA_TEXTO   = RGBColor(0x2A, 0x2A, 0x2A)
CINZA_SUB     = RGBColor(0x5A, 0x5A, 0x5A)
CINZA_CLARO   = RGBColor(0xF4, 0xF5, 0xF7)
LARANJA       = RGBColor(0xE8, 0x7A, 0x1E)
DOURADO       = RGBColor(0xC9, 0x96, 0x2C)

# ── Apresentação ────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W  = Inches(13.333)
SLIDE_H  = Inches(7.5)
MARGIN_L = Inches(0.8)
MARGIN_R = Inches(0.8)
CONTENT_W = SLIDE_W - MARGIN_L - MARGIN_R

BULLET  = "\u2022"   # •
DASH    = "\u2013"    # –
ARROW   = "\u2192"    # →
CHECK   = "\u2714"    # ✔
WARN    = "\u26A0"    # ⚠


# ── Helpers ─────────────────────────────────────────────────────────
def _bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def _box(slide, left, top, width, height):
    return slide.shapes.add_textbox(left, top, width, height)


def _font(run, size=18, bold=False, italic=False, color=CINZA_TEXTO, name="Calibri"):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = name


def _rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _thin_line(slide, top, color=LARANJA, thickness=Inches(0.04)):
    """Linha fina decorativa sob o cabeçalho."""
    _rect(slide, Inches(0), top, SLIDE_W, thickness, color)


# ── Tipos de slide ──────────────────────────────────────────────────
def slide_capa(titulo, subtitulo=""):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, AZUL_ESCURO)
    # Linha dourada decorativa
    _rect(s, Inches(3.5), Inches(3.1), Inches(6.333), Inches(0.04), DOURADO)
    # Título
    tb = _box(s, MARGIN_L, Inches(1.4), CONTENT_W, Inches(1.8))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = titulo
    _font(r, size=42, bold=True, color=BRANCO)
    # Subtítulo
    if subtitulo:
        tb2 = _box(s, MARGIN_L, Inches(3.4), CONTENT_W, Inches(2.0))
        tb2.text_frame.word_wrap = True
        p2 = tb2.text_frame.paragraphs[0]
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(12)
        for i, line in enumerate(subtitulo.split("\n")):
            px = p2 if i == 0 else tb2.text_frame.add_paragraph()
            px.alignment = PP_ALIGN.CENTER
            rx = px.add_run(); rx.text = line
            _font(rx, size=20, color=RGBColor(0xBB, 0xCC, 0xDD))


def slide_secao(titulo):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, AZUL_MEDIO)
    _rect(s, Inches(4.5), Inches(4.15), Inches(4.333), Inches(0.035), DOURADO)
    tb = _box(s, MARGIN_L, Inches(2.4), CONTENT_W, Inches(1.8))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = titulo
    _font(r, size=38, bold=True, color=BRANCO)


def _header_bar(slide, titulo):
    """Barra de cabeçalho + linha laranja."""
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.95), AZUL_ESCURO)
    _thin_line(slide, Inches(0.95))
    tb = _box(slide, MARGIN_L, Inches(0.15), CONTENT_W, Inches(0.7))
    tb.text_frame.word_wrap = True
    p = tb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = titulo
    _font(r, size=26, bold=True, color=BRANCO)


def slide_conteudo(titulo, items, subitems=None):
    """
    items: lista de strings (bullets principais).
    subitems: dict { indice: [sub1, sub2, ...] }
    """
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, BRANCO)
    _header_bar(s, titulo)

    tb = _box(s, MARGIN_L, Inches(1.25), CONTENT_W, Inches(5.8))
    tf = tb.text_frame
    tf.word_wrap = True
    if subitems is None:
        subitems = {}
    first = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(7)
        p.space_before = Pt(2)
        if item == "":
            r = p.add_run(); r.text = " "
            _font(r, size=10)
        else:
            r = p.add_run(); r.text = f" {BULLET}  {item}"
            _font(r, size=20, color=CINZA_TEXTO)
        # Sub-bullets
        if i in subitems:
            for sb in subitems[i]:
                ps = tf.add_paragraph()
                ps.space_after = Pt(3)
                ps.space_before = Pt(1)
                rs = ps.add_run(); rs.text = f"       {DASH}  {sb}"
                _font(rs, size=17, color=CINZA_SUB)
    return s


def slide_duas_colunas(titulo, tit_esq, items_esq, tit_dir, items_dir):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, BRANCO)
    _header_bar(s, titulo)

    col_w = Inches(5.6)
    gap = Inches(0.5)

    # Separador vertical sutil
    _rect(s, MARGIN_L + col_w + Inches(0.22), Inches(1.15),
           Inches(0.02), Inches(5.5), RGBColor(0xDD, 0xDD, 0xDD))

    for ci, (ct, cb) in enumerate([(tit_esq, items_esq), (tit_dir, items_dir)]):
        left = MARGIN_L + ci * (col_w + gap)
        # Título da coluna
        tb_t = _box(s, left, Inches(1.2), col_w, Inches(0.5))
        p_t = tb_t.text_frame.paragraphs[0]
        r_t = p_t.add_run(); r_t.text = ct
        _font(r_t, size=21, bold=True, color=AZUL_ACCENT)
        # Linha fina sob título de coluna
        _rect(s, left, Inches(1.68), Inches(3.0), Inches(0.02), AZUL_ACCENT)
        # Bullets
        tb_b = _box(s, left, Inches(1.85), col_w, Inches(5.0))
        tf = tb_b.text_frame
        tf.word_wrap = True
        for j, b in enumerate(cb):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.space_after = Pt(5)
            r = p.add_run(); r.text = f" {BULLET}  {b}"
            _font(r, size=17, color=CINZA_TEXTO)
    return s


def slide_tabela(titulo, headers, rows):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(s, BRANCO)
    _header_bar(s, titulo)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_h = Inches(min(5.5, 0.44 * n_rows))
    shape = s.shapes.add_table(n_rows, n_cols, MARGIN_L, Inches(1.25), CONTENT_W, tbl_h)
    tbl = shape.table

    # Cabeçalho
    for ci, h in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = AZUL_ESCURO
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for r in p.runs:
                _font(r, size=14, bold=True, color=BRANCO)

    # Linhas
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = CINZA_CLARO if ri % 2 == 0 else BRANCO
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for r in p.runs:
                    _font(r, size=13, color=CINZA_TEXTO)
    return s


# ====================================================================
#  CONTEÚDO DOS SLIDES
# ====================================================================

# ── 1. Capa ─────────────────────────────────────────────────────────
slide_capa(
    "Mapeamento de Análogos ao Terranimo",
    "No mundo e no Brasil — Foco em compactação por tráfego\n"
    "\n"
    "SoloCompactado IPT  |  30 / 03 / 2026"
)

# ── 2. Agenda ───────────────────────────────────────────────────────
slide_conteudo("Agenda", [
    "O que é o Terranimo e por que é referência",
    "Panorama internacional de ferramentas",
    "Quadro comparativo",
    "Análogos brasileiros: PredComp e COMPSOHMA",
    "Pacote soilphysics e PTFs brasileiras",
    "Grupos de pesquisa e bases de dados nacionais",
    "Lacunas críticas",
    "Próximos passos e proposta",
])

# ── 3. Terranimo ────────────────────────────────────────────────────
slide_conteudo("Terranimo: a referência internacional", [
    "Consórcio: Agroscope (CH), Aarhus (DK), SLU (SE), Wageningen (NL)",
    f"Modo Light: 4 entradas {ARROW} risco a 35 cm (semáforo verde / amarelo / vermelho)",
    "Modo Expert: solo por camadas (10 cm), textura completa, parâmetros mecânicos",
    "Banco de máquinas pré-populado + interface web amigável",
    "Integrado ao CLAAS CEMOS desde 2022 (decisão embarcada no trator)",
    "Nenhuma versão para solos tropicais / Brasil na página oficial",
], subitems={
    1: ["Carga na roda, pressão do pneu, teor de argila, sucção matricial"],
})

# ── 4. Seção Internacional ──────────────────────────────────────────
slide_secao("Panorama Internacional")

# ── 5. Modelos analíticos ───────────────────────────────────────────
slide_conteudo("Modelos analíticos (base científica)", [
    "SoilFlex (Keller et al. 2007) — motor do Terranimo",
    "TASC (Keller 2005) — sub-modelo de contato pneu-solo",
    "SOCOMO (van den Akker 2004) — precursor, Boussinesq + Froehlich",
    "COMPSOIL (O'Sullivan et al. 1999) — foco em variação de densidade",
    "PTFs Horn / Lebert (1991, 2003) — resistência do solo, base da política alemã",
], subitems={
    0: ["Tensão vertical via Boussinesq / Froehlich + contato super-elipse",
        "Implementado no pacote R soilphysics (função stressTraffic)"],
    3: ["Projeto EU SIDASS / ENVASSO; foco na magnitude (Δρ), não no risco"],
})

# ── 6. Ferramentas aplicadas, MEF e fronteira ──────────────────────
slide_conteudo("Ferramentas aplicadas, MEF e fronteira", [
    "REPRO (Alemanha) — sustentabilidade de fazenda; compactação é um módulo entre vários",
    "FRIDA (Alemanha) — diagnóstico qualitativo de risco por rodado",
    "Softsoil / PLAXIS (Bentley) — MEF comercial, muito detalhado, impraticável p/ campo",
    "Abordagens ANN / ML (2020–2025) — data-driven, sem ferramenta operacional ainda",
    "Soil2Cover (Keller et al. 2025) — FRONTEIRA: otimização de rota + SoilFlex",
], subitems={
    4: ["Não só avalia risco, mas otimiza a rota que minimiza compactação acumulada",
        "Testado em 1 000+ campos. Publicação: Precision Agriculture 2025"],
})

# ── 7. Quadro comparativo ──────────────────────────────────────────
slide_tabela("Quadro comparativo internacional", [
    "Ferramenta", "País", "Web?", "Status", "Diferencial"
], [
    ["Terranimo",        "CH/DK/SE", "Sim",   "Ativo",     "Referência: motor + UX + banco de máquinas"],
    ["SoilFlex",         "CH/SE",    "Via R",  "Ativo",     "Base científica do Terranimo"],
    ["TASC",             "CH/SE",    "Não",   "Ativo",     "Sub-modelo de contato pneu"],
    ["SOCOMO",           "NL",       "Não",   "Legado",    "Precursor analítico"],
    ["COMPSOIL",         "IE/EU",    "Não",   "Acadêmico", "Foco em Δρ"],
    ["REPRO",            "DE",       "Parcial","Incerto",   "Módulo em ferramenta de fazenda"],
    ["Softsoil / PLAXIS","NL",       "Não",   "Comercial", "MEF detalhado"],
    ["Soil2Cover",       "CH/SE",    "Não",   "Pesquisa",  "Otimização de rota + SoilFlex"],
    ["PredComp",         "BR",       "Sim",   "Ativo",     "SoilFlex em Shiny"],
    ["COMPSOHMA",        "BR",       "Sim",   "Beta",      "Light brasileiro, foco cana"],
    ["Horn PTFs",        "DE",       "Não",   "Política",  "Resistência p/ regulação"],
    ["ANN / ML",         "Vários",   "Não",   "Pesquisa",  "Data-driven, sem ferramenta pronta"],
])

# ── 8. Seção Brasil ────────────────────────────────────────────────
slide_secao("Análogos Brasileiros")

# ── 9. PredComp vs COMPSOHMA ───────────────────────────────────────
slide_duas_colunas(
    "PredComp vs COMPSOHMA",
    "PredComp (acadêmico)",
    [
        "Base: soilphysics / SoilFlex",
        "5 abas: tensão, pré-adensamento,\n  prop. compressivas, risco, predição",
        "Entradas detalhadas: pressão pneu,\n  carga, N, λ, κ, ρ …",
        "Exige conhecimento de parâmetros\n  mecânicos do solo",
        "Ativo no CRAN, registro INPI",
        "Financiamento FAPESP + CNPq",
    ],
    "COMPSOHMA (aplicado, beta)",
    [
        "Estilo Terranimo Light",
        "Risco até 1 m (4 camadas)",
        "Cenários pré-configurados (cana)",
        "Entradas simples: máquina, argila\n  por camada, umidade qualitativa",
        "Saída: perfil verde / amarelo / vermelho",
        "Beta, registro INPI",
        "Foco atual: cana-de-açúcar",
    ],
)

# ── 10. soilphysics ────────────────────────────────────────────────
slide_tabela("Pacote soilphysics — funções de compactação", [
    "Função", "O que calcula"
], [
    ["stressTraffic()",          "Área de contato + tensão + propagação (Keller 2005 + Froehlich)"],
    ["soilDeformation()",        "Variação de ρ após tensão (O'Sullivan & Robertson 1996)"],
    ["soilStrength()",           "σp via PTFs (Lebert / Horn, Imhoff, Severiano etc.)"],
    ["soilStrength2()",          "PTFs adicionais para σp"],
    ["compressive_properties()", "N, λn, κ via PTFs (O'Sullivan, Defossez, de Lima)"],
])

# ── 11. PTFs brasileiras ───────────────────────────────────────────
slide_conteudo("PTFs calibradas com solos brasileiros (já no pacote)", [
    "Severiano et al. (2013): σp para Latossolos, argila 152–521 g/kg",
    "Imhoff et al. (2004): σp para Hapludox (ρ, CC, w)",
    "de Lima et al. (2018): N, λ, κ para Franco-arenoso e Franco-argiloarenoso",
    "de Lima et al. (2020b): PTFs para camada arável e pé-de-grade",
    "",
    f"{CHECK}  Cobertura boa para Latossolos / Oxissolos",
    f"{WARN}  Lacuna: Argissolos, Nitossolos, Cambissolos, Neossolos",
], subitems={
    0: ["σp(152) = 129 h⁰·¹⁵  …  σp(521) = 63 h⁰·¹⁵"],
})

# ── 12. Grupos de pesquisa ─────────────────────────────────────────
slide_conteudo("Grupos de pesquisa (Brasil)", [
    "ESALQ / USP (Piracicaba) — HUB PRINCIPAL",
    "IF Goiano (Urutaí) — A. R. da Silva, co-desenvolvedor soilphysics",
    "UFLA (Lavras) — Severiano et al., PTFs para Latossolos",
    "UEM / UEPG (Paraná) — Tormena, Giarola: colaboração com Keller (SoilFlex-LLWR)",
    "UFRPE (Recife) — ensaios de rodagem com solos brasileiros",
    "EMBRAPA Solos — BDSolos, SiBCS, app SoloClass (sem ferramenta de compactação)",
], subitems={
    0: ["R. P. de Lima (soilphysics, PredComp, COMPSOHMA)",
        "M. R. Cherubin (COMPSOHMA, qualidade do solo em cana)",
        "A. P. da Silva (co-autor soilphysics, pesquisador sênior)",
        "R. B. Menillo (probe automatizada de ρ, patente INPI)"],
})

# ── 13. Bases de dados ─────────────────────────────────────────────
slide_tabela("Bases de dados e recursos brasileiros", [
    "Recurso", "Instituição", "Relevância para protótipo nacional"
], [
    ["BDSolos",          "EMBRAPA",        f"Milhares de perfis (granulometria, ρ, MO) {ARROW} parâmetros default"],
    ["HYBRAS",           "ANA / INPE",     f"Curvas de retenção {ARROW} estimar sucção / umidade"],
    ["SiBCS",            "EMBRAPA",        f"Classificação {ARROW} mapear vulnerabilidade à compactação"],
    ["MapBiomas Solos",  "Rede MapBiomas", "Carbono orgânico espacializado"],
    ["SoloClass (app)",  "EMBRAPA",        "Classificação de campo no celular"],
])

# ── 14. Lacunas ────────────────────────────────────────────────────
slide_conteudo("Lacunas críticas para um Terranimo brasileiro", [
    "Sem ferramenta agnóstica de cultura (COMPSOHMA só cana)",
    f"Sem integração espacial (GPS / RTK {ARROW} mapa de risco por trilha)",
    "Sem capacidade embarcada (nada como CLAAS CEMOS + Terranimo)",
    "Cobertura incompleta de PTFs (faltam Argissolos, Nitossolos etc.)",
    "Sem conexão automática com BDSolos / SiBCS (entrada toda manual)",
    "Sem framework regulatório (Brasil não tem equivalente ao BBodSchG alemão)",
])

# ── 15. Próximos passos ────────────────────────────────────────────
slide_conteudo("Próximos passos e proposta", [
    "Benchmark técnico: PredComp / soilphysics (SoilFlex em R, PTFs prontas)",
    "Benchmark de interface: COMPSOHMA + Terranimo Light",
    "Fronteira: Soil2Cover (otimização de rota + compactação acumulada)",
    "",
    "Tropicalizar o modelo:",
    "",
    "Integração-alvo:",
], subitems={
    4: ["Granulometria, Atterberg, índice L, índice de vazios",
        "Grau de saturação, histórico de umedecimento / secagem",
        "2–4 classes de solo piloto → calibrar σcrit por camada"],
    6: [f"SiBCS / BDSolos {ARROW} PTFs automáticas {ARROW} modelo de risco por camada",
        "Saída operacional: semáforo + recomendação de ajuste"],
})

# ── 16. Encerramento ───────────────────────────────────────────────
slide_capa(
    "Obrigado",
    "SoloCompactado IPT\n"
    "Documentação completa:\n"
    "docs/mapeamento_analogos_terranimo_brasil_2026-03-30.pdf"
)

# ── Salvar ──────────────────────────────────────────────────────────
out = "/home/ggrl/projetos/ipt/Civil/Geotecnia/SoloCompactado/docs/apresentacao_analogos_terranimo_2026-03-30.pptx"
prs.save(out)
print(f"Salvo: {out}")
print(f"Slides: {len(prs.slides)}")
